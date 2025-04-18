import random 
import os
import math
from pptx import Presentation
from pptx.util import Inches
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.util import Cm, Pt
import freepik 

class Slide:
    def __init__(self, header, body, type):
        self.header = header    
        self.body = body
        self.type = type


# For Fz Rubik font:  [max_char, font_size, width, height]
ratio = [27, 33, 6739779, 507831] # Subhead
ratio2 = [55, 16, 6739779, 507831] # Subcont

def create_ppt(text_file, design_number, ppt_name):
    slide_count = 0
    header = ""
    body = ""
    type = ""
    #--------------------------------#
    # Make a list of images for illusion
    gen_image_list = []
    for root, dirs, files in os.walk('img'):
        for filenames in files:
            if (filenames.endswith(('.png', '.jpg', '.jpeg', '.tiff', '.bmp', '.gif'))):
                gen_image_list.append(filenames)

    #--------------------------------#
    # Create a list for storing slide text
    slide_text = []
    with open(text_file, 'r', encoding='utf-8') as f:
        for line_num, line in enumerate(f):
            if line.startswith('#Title:'):
                header = line.replace('#Title:', '').strip()
                type = "Title"
                slide = Slide(header, body, type)
                slide_text.append(slide)

                header = ""
                body = ""
                type = ""
                continue

            elif line.startswith('#Slide:'):
                if (slide_count > 0):
                    # Type of content
                    body = "\n".join(body.splitlines()[1:]) if body.startswith('\n') else body
                    if (slide_count == 1):
                        type = "Table"
                    elif (body[0].isdigit()):
                        type = "Content2"
                        if (header.lower() == "summary"):
                            type = "Summary"
                
                slide = Slide(header, body, type)
                slide_text.append(slide)

                slide_count += 1
                header = ""
                body = ""
                type = ""
                continue

            elif line.startswith('#Header:'):
                header = line.replace('#Header:', '').strip()
                type = "Content"
                continue

            elif line.startswith('#Content:'):
                body = line.replace('#Content:', '').strip()
                next_line = f.readline().strip()
                while next_line and not next_line.startswith('#'):
                    if (next_line):
                        body += '\n' + next_line
                    next_line = f.readline().strip()
                continue
            
    
    #-----------------------------#
    # Make slide for each part
    # Slide 0: Text 12, picture 1 -- Title
    # Slide 1: Max 6 Heading
    # Slide 2: 2 Content per slide. Each contains: Title - Heading - Content - Picture
    # Slide 3: 1 Content per slide. Each contains: Title - Heading - Content - Picture
    # Slide 4: 1 Content per slide, but have at least 2 subheadings. 
    # Each contains: Title - Heading 
    # Subheading1 - Content - Picture
    prs = Presentation(f"Designs/Design-{design_number}.pptx")
    content_count = 0

    #Create list for placeholder indexs
    placeholder_list = []
    index = 0
    remained_slot = 0
    content_slide_layout = [2, 3]
    content_slide_count = 0
    slide_layout_index = 2
    first_time = True
    hasSubheading = False

    for part in slide_text:
        if (type == "Content"):
            content_slide_count += 1

    # for part in slide_text:
    #     print(part.header, "\n", part.body, "\n", part.type)
    #     print("----------")

    rad = [6879185, 246221, 58, 16]
    #width, max_character_per_line

    for part in slide_text:
        header = part.header
        body = part.body
        type = part.type
        # Title slide
        if (type == "Title"):
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            title = slide.shapes.placeholders[12]
            title.text = header

        if (type == "Table"):
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            #Create list for placeholder indexs
            text_placeholder_list = []
            for shape in slide.placeholders:
                id = shape.placeholder_format.idx
                name = shape.name
                if (name.startswith('Text')):
                    text_placeholder_list.append(id)
            #Header of current slide
            slide.placeholders[text_placeholder_list[0]].text = header
            #Separate to each part in table of content
            lines = body.split("\n")
            line_count = 0
            list_idx = 1
            for line in lines:
                if (line):
                    stripped_line = line.split(". ", 1)[1]
                    line_count += 1
                    if (list_idx + 1 < len(text_placeholder_list)):
                        slide.placeholders[text_placeholder_list[list_idx]].text = "0" + str(line_count) #Number
                        slide.placeholders[text_placeholder_list[list_idx + 1]].text = stripped_line #Text
                        list_idx += 2

        if (type == "Content"):
            if (remained_slot == 0 or first_time):
                if (content_slide_count > 1):
                    slide_layout_index = random.choice(content_slide_layout)
                else: 
                    slide_layout_index = 3
                if (slide_layout_index == 2): remained_slot = 2
                else: remained_slot = 1
                slide = prs.slides.add_slide(prs.slide_layouts[slide_layout_index])
                placeholder_list = []
                index = 0
                first_time = False
                for shape in slide.placeholders:
                    id = shape.placeholder_format.idx
                    name = shape.name
                    placeholder_list.append(id)

            content_count += 1
            # For layout 2 only
            if (slide_layout_index == 2):
                if (remained_slot == 1):
                    index += 4

            slide.placeholders[placeholder_list[index]].text = "0" + str(content_count)

            base_w = rad[0]
            base_h = rad[1]
            base_max_char = rad[2]
            base_font = rad[3]

            font_size = []
            #Change font size of heading
            w = slide.placeholders[placeholder_list[index + 1]].width
            h = slide.placeholders[placeholder_list[index + 1]].height
            current_length = len(header)
            this_ratio = (w * h) / (base_w * base_h)
            calculated_font_size = base_font * math.sqrt(this_ratio) * math.sqrt(base_max_char/ current_length)
            slide.placeholders[placeholder_list[index + 1]].text = header
            title_para = slide.shapes.placeholders[placeholder_list[index + 1]].text_frame.paragraphs[0]
            title_para.font.size = Pt(calculated_font_size)
            font_size.append(calculated_font_size)
            #Change font size of content
            w = slide.placeholders[placeholder_list[index + 2]].width
            h = slide.placeholders[placeholder_list[index + 2]].height
            current_length = len(body)
            this_ratio = (w * h) / (base_w * base_h)
            calculated_font_size = base_font * math.sqrt(this_ratio) * math.sqrt(base_max_char/ current_length)
            slide.placeholders[placeholder_list[index + 2]].text = body
            title_para = slide.shapes.placeholders[placeholder_list[index + 2]].text_frame.paragraphs[0]
            title_para.font.size = Pt(calculated_font_size)
            font_size.append(calculated_font_size)

            #Handle image placeholder
            w = slide.placeholders[placeholder_list[index + 3]].width
            h = slide.placeholders[placeholder_list[index + 3]].height
            choosed_picture = freepik.process(body, w, h)
            slide.placeholders[placeholder_list[index + 3]].insert_picture("img/" + choosed_picture)
            remained_slot -= 1
        
        if (type == "Content2"):
            placeholder_list = []
            remained_slot = 1

            content_with_subhead = [4, 6]
            content_index = random.choice(content_with_subhead)
            slide = prs.slides.add_slide(prs.slide_layouts[content_index])

            print(content_index)
            slide_count += 1

            lines = body.split("\n")
            #Create list for image placeholder indexs
            image_placeholder_list = []
            #Create list for text placeholder indexs
            text_placeholder_list = []

            title_index = 0
            for shape in slide.placeholders:
                id = shape.placeholder_format.idx
                name = shape.name
                if (name.startswith('Title')):
                    title_index = id
                if (name.startswith('Text')):
                    text_placeholder_list.append(id)
                if (name.startswith('Picture')):
                    image_placeholder_list.append(id)

            text_list_idx = 0
            img_list_idx = 0
            content_count += 1
            #Title  
            slide.placeholders[title_index].text = "0" + str(content_count)
            #Heading
            slide.placeholders[text_placeholder_list[0]].text = header
            text_placeholder_list.pop(0)
            #Subheading
            font_size = []
            title_para = slide.shapes.placeholders[text_placeholder_list[0]].text_frame.paragraphs[0]
            font_size.append(31)
            for line in lines:
                if (line):
                    content = line.split(". ")
                    subhead = content[1]
                    subcontent = content[2]

                    if (text_list_idx + 1 < len(text_placeholder_list)):
                        cur_len = len(subhead)
                        base_w = rad[0]
                        base_h = rad[1]
                        base_max_char = rad[2]
                        base_font = rad[3]

                        #Change font size of heading
                        w = slide.placeholders[text_placeholder_list[text_list_idx + 1]].width
                        h = slide.placeholders[text_placeholder_list[text_list_idx + 1]].height
                        current_length = len(subcontent)
                        this_ratio = (w * h) / (base_w * base_h)
                        calculated_font_size = base_font * math.sqrt(this_ratio) * math.sqrt(base_max_char/ current_length)
                        font_size.append(int(calculated_font_size))
                        slide.placeholders[text_placeholder_list[text_list_idx]].text = subhead
                        slide.placeholders[text_placeholder_list[text_list_idx + 1]].text = subcontent
                        text_list_idx += 2

                    if (img_list_idx < len(image_placeholder_list)):
                        w = slide.placeholders[image_placeholder_list[img_list_idx]].width
                        h = slide.placeholders[image_placeholder_list[img_list_idx]].height
                        choosed_picture = freepik.process(subcontent, w, h)
                        slide.placeholders[image_placeholder_list[img_list_idx]].insert_picture("img/" + choosed_picture)
                        img_list_idx += 1
            
            remained_slot -= 1
            min_font_size = min(font_size)
            
            print("fz: ", font_size)
            text_list_idx = 0
            while text_list_idx + 1 < len(text_placeholder_list):      
                title_para = slide.shapes.placeholders[text_placeholder_list[text_list_idx + 1]].text_frame.paragraphs[0]
                title_para.font.size = Pt(min_font_size)
                text_list_idx += 2
            # print(image_placeholder_list, "hihi", text_placeholder_list)

        #Summary slide
        if (type == "Summary"):
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            title = slide.shapes.title
            title.text = header
            #Create list for image placeholder indexs
            image_placeholder_list = []
            #Create list for text placeholder indexs
            text_placeholder_list = []
            for shape in slide.placeholders:
                id = shape.placeholder_format.idx
                name = shape.name
                if (name.startswith('Text')):
                    text_placeholder_list.append(id)
                if (name.startswith('Picture')):
                    image_placeholder_list.append(id)

            for index in text_placeholder_list:
                slide.placeholders[index].text = str(index)
            text_list_idx = 1
            img_list_idx = 0
            content_count += 1

            #Title  
            slide.placeholders[0].text = "0" + str(content_count)
            #Heading
            slide.placeholders[text_placeholder_list[0]].text = header
            #Subheading
            lines = body.split("\n")

            min_size = []
            for line in lines:
                if (line):
                    content = line.split(". ")
                    subhead = content[0]
                    subcontent = content[1]

                    base_w = rad[0]
                    base_h = rad[1]
                    base_max_char = rad[2]
                    base_font = rad[3]

                    w = slide.placeholders[text_placeholder_list[text_list_idx]].width
                    h = slide.placeholders[text_placeholder_list[text_list_idx]].height
                    current_length = len(subcontent)
                    this_ratio = (w * h) / (base_w * base_h)
                    calculated_font_size = base_font * math.sqrt(this_ratio) * math.sqrt(base_max_char/ current_length)
                    min_size.append(int(calculated_font_size))
                    slide.placeholders[text_placeholder_list[text_list_idx]].text = subcontent
                    # title_para = slide.shapes.placeholders[text_placeholder_list[text_list_idx]].text_frame.paragraphs[0]
                    # title_para.font.size = Pt(calculated_font_size)
                    text_list_idx += 1

            min_font_size = min(min_size)
            for text_list_idx in range(1, len(text_placeholder_list)):
                slide.shapes.placeholders[index].text_frame.paragraphs[0].font.size = Pt(min_font_size) 
                title_para = slide.shapes.placeholders[text_placeholder_list[text_list_idx]].text_frame.paragraphs[0]
                title_para.font.size = Pt(min_font_size)  
        
            w = slide.placeholders[image_placeholder_list[img_list_idx]].width
            h = slide.placeholders[image_placeholder_list[img_list_idx]].height    
            choosed_picture = freepik.process(body, w, h)
            slide.placeholders[image_placeholder_list[img_list_idx]].insert_picture("img/" + choosed_picture)

    prs.save(f'GeneratedPresentations/{ppt_name}.pptx')

filename = "robot"
pptlink = create_ppt(f'Cache/{filename}.txt', 9, filename)
