from flask import Flask, render_template, request, send_from_directory, redirect, jsonify, send_file
from flask_limiter import Limiter
from flask_limiter.util import get_remote_address
import os
import openai
import collections.abc
from pptx import Presentation
from pptx.util import Inches
import random 
import re

client = openai.OpenAI(api_key='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')

app = Flask(__name__)

limiter = Limiter(
    app,
    # default_limits=["10 per day"], #This is the rate limit, you can remove it if you want
)
#-----------------------------------------------#
import prompt
def create_ppt_text(Input):
    response = client.chat.completions.create(
        model="gpt-4o",
        messages=[
            {"role": "system", "content": (prompt.Content_Prompt)},
            {"role": "user", "content": ("The user wants a presentation about " + Input)}
        ],
        temperature=0.5,
    )
    return response.choices[0].message.content

def extract_keyword(Content):
    response = client.chat.completions.create(
        model="gpt-4o",
        messages=[
            {"role": "system", "content": (prompt.Extract_Prompt)},
            {"role": "user", "content": ("Extract some keywords from following text: " + Content)}
        ],
        temperature=0.5,
    )
    answer = response.choices[0].message.content
    print(answer)
    return answer
#-----------------------------------------------#

class Slide:
    def __init__(self, header, body, type):
        self.header = header    
        self.body = body
        self.type = type

def create_ppt(text_file, design_number, ppt_name):
    prs = Presentation(f"Designs/Design-{design_number}.pptx")
    slide_count = 0
    header = ""
    content = ""
    last_slide_layout_index = -1
    firsttime = True

    #--------------------------------#
    # Make a list of images for illusion
    gen_image_list = []
    for root, dirs, files in os.walk('img'):
        for filenames in files:
            if (filenames.endswith(('.png', '.jpg', '.jpeg', '.tiff', '.bmp', '.gif'))):
                gen_image_list.append(filenames)

    #--------------------------------#
    # Create a list for storing slide text

    slide_text = []
    with open(text_file, 'r', encoding='utf-8') as f:
        for line_num, line in enumerate(f):
            if line.startswith('#Title:'):
                header = line.replace('#Title:', '').strip()
                slide = prs.slides.add_slide(prs.slide_layouts[0])
                title = slide.shapes.title
                title.text = header
                body_shape = slide.shapes.placeholders[1]

                header = line.replace('#Title:', '').strip()
                body = ""
                type = "Title"
                slide = Slide(header, body, type)
                slide_text.append(slide)

                continue
            elif line.startswith('#Slide:'):

                if slide_count > 0:
                    slide = prs.slides.add_slide(prs.slide_layouts[slide_layout_index])

                    #Create list for image placeholder indexs
                    image_placeholder_list = []
                    #Create list for placeholder indexs
                    text_placeholder_list = []

                    for shape in slide.placeholders:
                        id = shape.placeholder_format.idx
                        name = shape.name
                        if (name.startswith('Text')):
                            text_placeholder_list.append(id)
                        if (name.startswith('Picture')):
                            image_placeholder_list.append(id)

                    print(len(image_placeholder_list), len(text_placeholder_list), slide_layout_index, slide_placeholder_index)
                    print(content)
                    # Adding title
                    title = slide.shapes.title
                    title.text = header
                    # Slide 1 only:
                    if (slide_layout_index == 1):
                        body = slide.placeholders[slide_placeholder_index]
                        body.text = content
                    # Adding image if needed, for Slide 8              
                    elif (image_placeholder_list and text_placeholder_list): # Image insert, except Slide 1
                        body = slide.placeholders[text_placeholder_list[0]]
                        body.text = content
                        selected_pic = random.choice(gen_image_list)
                        slide.placeholders[image_placeholder_list[0]].insert_picture("img/" + selected_pic)

                content = "" 
                slide_count += 1
                slide_layout_index = last_slide_layout_index
                layout_indices = [1, 2] 
                while slide_layout_index == last_slide_layout_index:
                    if firsttime == True:
                        slide_layout_index = 1
                        slide_placeholder_index = 1
                        firsttime = False
                        break
                    slide_layout_index = random.choice(layout_indices) # Select random slide index
                    if slide_layout_index == 7:
                        slide_placeholder_index = 2 # text
                    else:
                        slide_placeholder_index = 1 # picture
                last_slide_layout_index = slide_layout_index
                continue

            elif line.startswith('#Header:'):
                header = line.replace('#Header:', '').strip()
                continue

            elif line.startswith('#Content:'):
                content = line.replace('#Content:', '').strip()
                next_line = f.readline().strip()
                while next_line and not next_line.startswith('#'):
                    content += '\n' + next_line
                    next_line = f.readline().strip()
                continue
    
    prs.save(f'GeneratedPresentations/{ppt_name}.pptx')
    file_path = f"GeneratedPresentations/{ppt_name}.pptx"
    return f"{request.host_url}{file_path}"

#----------------------------------------------------#

@app.route('/upload', methods=['POST'])
def upload_file():
    if 'file' not in request.files:
        return "No file part"
    file = request.files['file']
    if file.filename == '':
        return "No selected file"
    # Save the file or process its content
    file_content = file.read().decode('utf-8')
    # Analyze the file content (e.g., send it to ChatGPT)
    return f"File uploaded and analyzed successfully! Content: {file_content[:500]}"  # Example: show the first 500 characters

@app.route('/GeneratedPresentations/<path:path>')
def send_generated_image(path):
    return send_file(f'GeneratedPresentations/{path}', as_attachment=True)
    
@app.route("/powerpoint")
def powerpoint():
    return render_template("powerpoint.html", charset="utf-8")
    
@app.route("/")
def home():
    return render_template("powerpoint.html", charset="utf-8")

@app.route("/get")
@limiter.limit("10 per day, key_func=get_remote_address")
def get_bot_response():
    user_text = request.args.get("msg")
    last_char = user_text[-1]
    input_string = user_text
    input_string = re.sub(r'[^\w\s.\-\(\)]', '', input_string)
    input_string = input_string.replace("\n", "")
    number = 1

    if last_char.isdigit():
        number = int(last_char)
        input_string = user_text[:-2]
        print("Design Number:", number, "selected.")
    else:
        print("No design specified, using default design...")
        
    if number > 8:
        number = 1
        print("Unavailable design, using default design...")
    elif number == 0:
        number = 1
        print("Unavailable design, using default design...")

    # Generate a filename using OpenAI API
    filename_prompt = f"Generate a short, descriptive filename based on the following input: \"{input_string}\". Answer just with the short filename, no other explainment."
    filename_response = client.chat.completions.create(
        model="gpt-4o",
        messages=[
            {"role": "system", "content": filename_prompt},
        ],
        temperature=0.5,
        max_tokens=30,
    )
    filename = filename_response.choices[0].message.content.strip().replace(" ", "_")

    with open(f'Cache/{filename}.txt', 'w', encoding='utf-8') as f:
        f.write(create_ppt_text(input_string))

    pptlink = create_ppt(f'Cache/{filename}.txt', number, filename)
    return str(pptlink)

if __name__ == '__main__':
    # debug=True in the run() parameter if you want to debug
    app.run()